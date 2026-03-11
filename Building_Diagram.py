import os
import shutil
import time
from io import BytesIO
from pptx import Presentation

# Absolute path to the project source directory — used for reading templates
_SRC_DIR = os.path.dirname(os.path.abspath(__file__))


def rearrange_folders():
    """Moves any existing ×JUNC× Desktop folder to a timestamped archive folder,
    then creates a fresh ×JUNC× folder. Only used in standalone (non-Streamlit) mode."""
    desktop = os.path.join(os.environ['USERPROFILE'], 'Desktop')
    old_path = os.path.join(desktop, "×JUNC×")
    if os.path.exists(old_path):
        old_id_path = os.path.join(old_path, "×ID×.png")
        if len(os.listdir(old_path)) == 0:
            os.rmdir(old_path)
            os.makedirs(old_path)
        else:
            if os.path.exists(old_id_path):
                path_created = os.path.getmtime(old_id_path)
            else:
                path_created = os.path.getmtime(old_path)
            year, month, day, hour, minute, second = time.localtime(path_created)[:-3]
            folder = "×JUNC× " + str(" %02d⌟%02d⌟%d %02d∶%02d" % (day, month, year, hour, minute))
            new_path = os.path.join(desktop, folder)
            if not os.path.exists(new_path):
                os.makedirs(new_path)
            else:
                i = 1
                new_folder = False
                while not new_folder:
                    if not os.path.exists(new_path + "[%s]" % i):
                        new_path = new_path + "[%s]" % i
                        os.makedirs(new_path)
                        new_folder = True
                    else:
                        i += 1
            for file_name in os.listdir(old_path):
                shutil.move(os.path.join(old_path, file_name), new_path)
    else:
        os.makedirs(old_path)


def create_new_diagram_template_file():
    """Creates a working copy of the Diagram template in the current directory."""
    src = os.path.join(_SRC_DIR, "Diagram_template", "Diagram_templates.pptx")
    prs = Presentation(src)
    prs.save("Diagram_new_template.pptx")


def del_slides(pres, chosen_type):
    """Deletes all slides from the template that don't match the chosen junction type."""
    i = len(pres.slides)
    while i > 0:
        if i != chosen_type:
            rId = pres.slides._sldIdLst[i - 1].rId
            pres.part.drop_rel(rId)
            del pres.slides._sldIdLst[i - 1]
        i -= 1
    pres.save("Del_Diagram.pptx")


def get_pptx_bytes(pres):
    """Serialize a Presentation object to bytes (for Streamlit downloads)."""
    buf = BytesIO()
    pres.save(buf)
    buf.seek(0)
    return buf.getvalue()


def delete_temp_diagram_pres():
    """Deletes all temporary PPTX files created during diagram generation."""
    prs_to_del = [
        "Diagram_new_template.pptx", "Del_Diagram.pptx", "Street_Diagram.pptx",
        "Morn_Diagram.pptx", "Eve_Diagram.pptx", "Dirc_Diagram.pptx",
    ]
    for temp_prs in prs_to_del:
        if os.path.exists(temp_prs):
            os.remove(temp_prs)
