import os
from pptx import Presentation

# Absolute path to the project source directory — used for reading templates
_SRC_DIR = os.path.dirname(os.path.abspath(__file__))


def create_new_table_templates_file():
    """Creates a working copy of the Table template in the current directory."""
    src = os.path.join(_SRC_DIR, "Table_template", "Table_templates.pptx")
    prs_tbl = Presentation(src)
    prs_tbl.save("Table_new_template.pptx")


def del_slides_table(pres, chosen_type):
    """Deletes all slides from the template that don't match the chosen table type."""
    i = len(pres.slides)
    while i > 0:
        if i != chosen_type:
            rId = pres.slides._sldIdLst[i - 1].rId
            pres.part.drop_rel(rId)
            del pres.slides._sldIdLst[i - 1]
        i -= 1
    pres.save("Del_Table.pptx")


def delete_temp_table_pres():
    """Deletes all temporary PPTX files created during table generation."""
    prs_to_del = [
        "Table_new_template.pptx", "Del_Table.pptx", "Vol_Table.pptx",
        "Info_Table.pptx", "Dirc_Table.pptx",
    ]
    for temp_prs in prs_to_del:
        if os.path.exists(temp_prs):
            os.remove(temp_prs)
