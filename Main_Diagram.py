# Version date 8/6/2021

import os
from io import BytesIO
import Phaser
from Directions import *
from General_Info import *
from LRT_Info import *
from Phaser_Output import *
from Building_Diagram import *
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.util import Pt

_JUNC_DIR = os.path.dirname(os.path.abspath(__file__))


class Diagram:
    """
    A class used to represent all the info about the junction that is needed for creating the Diagram presentation file.
    """

    def __init__(self, output_phaser_list, xlprop):
        """ The constructor of the Diagram class,called when a new instance of a class is created.
         To initialize, it needs the output of Phaser and the info from the excel"""
        self.__North = Direction(
            "North")  # A property representing the north direction. Initialized with the name "North"
        self.__South = Direction(
            "South")  # A property representing the south direction. Initialized with the name "South"
        self.__East = Direction("East")  # A property representing the east direction. Initialized with the name "East"
        self.__West = Direction("West")  # A property representing the west direction. Initialized with the name "West"
        self.__GenInfo = General_Info()  # A property representing general info about the junction.
        self.__LRTInfo = LRT_Info()  # A property representing LRT info about the junction.
        self.__InfoFromPhaser = PhsrOutput(
            output_phaser_list)  # A class that holds the output from Phaser. That is the data the will be pushed to
        # the different class in Diagram and Table.
        self.__ExcelProperties = xlprop  # A list with info about the creator of the volume_calculator excel file. It is
        # used in the ID file
        self.WARNINGS = []  # Collects non-fatal warnings to surface in the UI

    @property
    def NO(self):
        """Get the north info"""
        return self.__North

    @NO.setter
    def NO(self, value):
        """Set the north info"""
        self.__North = value

    @property
    def SO(self):
        """Get the south info"""
        return self.__South

    @SO.setter
    def SO(self, value):
        """Set the south info"""
        self.__South = value

    @property
    def EA(self):
        """Get the east info"""
        return self.__East

    @EA.setter
    def EA(self, value):
        """Set the east info"""
        self.__East = value

    @property
    def WE(self):
        """Get the west info"""
        return self.__West

    @WE.setter
    def WE(self, value):
        """Set the north info"""
        self.__West = value

    @property
    def G_INF(self):
        """Get the general info"""
        return self.__GenInfo

    @G_INF.setter
    def G_INF(self, value):
        """Set the general info"""
        self.__GenInfo = value

    @property
    def LRT_INF(self):
        """Get the LRT info"""
        return self.__LRTInfo

    @LRT_INF.setter
    def LRT_INF(self, value):
        """Set the LRT info"""
        self.__LRTInfo = value

    @property
    def phsr_lst(self):
        """Get the phaser list info"""
        return self.__InfoFromPhaser

    @phsr_lst.setter
    def phsr_lst(self, value):
        """Set the phaser list info"""
        self.__InfoFromPhaser = value

    @property
    def xlprop(self):
        """Get the excel properties of the current JUNC"""
        return self.__ExcelProperties

    @xlprop.setter
    def xlprop(self, value):
        """Set the excel properties of the current JUNC"""
        self.__ExcelProperties = value

    def push_arr(self):
        """the method uses the output arrows of Phaser to push them into the right subclass of each direction,
        divided to regular arrows and public transport arrows. """
        arr_list = [self.phsr_lst.ARROW_REG, self.phsr_lst.ARROW_PT]
        orig_lanes = ["R", "TR", "T", "TL", "L", "A", "RL"]
        directions = [self.NO.LAN, self.SO.LAN, self.EA.LAN, self.WE.LAN]

        for direc in directions:
            for lan in orig_lanes:
                if arr_list[0]:
                    cur_arrow_input = [arr_list[0][0], arr_list[1][0]]
                    if lan == "R":
                        SR_cur_arrow_input = [0, 0]
                        if arr_list[0][0] == 9:
                            SR_cur_arrow_input[0] = 1
                            cur_arrow_input[0] = 0
                        if arr_list[1][0] == 9:
                            SR_cur_arrow_input[1] = 1
                            cur_arrow_input[1] = 0
                        if sum(SR_cur_arrow_input) > 0:
                            setattr(direc, "SR", SR_cur_arrow_input)
                    setattr(direc, lan, cur_arrow_input)
                    arr_list[0].pop(0)
                    arr_list[1].pop(0)

    def push_vol(self):
        """the method uses the output volumes of Phaser to push them into the right subclass of each direction,
        divided to morning and evening"""
        vol_list = [self.phsr_lst.MOR_VOL, self.phsr_lst.EVE_VOL]
        directions_mor = [self.NO.MOR, self.SO.MOR, self.EA.MOR, self.WE.MOR]
        directions_eve = [self.NO.EVE, self.SO.EVE, self.EA.EVE, self.WE.EVE]
        count = -1

        for vol in vol_list:
            count += 1
            if count == 0:
                directions = directions_mor
            else:
                directions = directions_eve
            for direc in directions:
                routes = ["R", "T", "L"]
                for rou in routes:
                    if vol:
                        value_to_push = int(vol[0])
                        setattr(direc, rou, value_to_push)
                        vol.pop(0)

    def push_general_info(self):
        """the method uses the output general information of Phaser to push it into G_INF subclass and to each
        matching property in that subclass. For specific info that related to the LRT, it pushes it to LRT_INF """
        phaser_gen_info_list = self.phsr_lst.GEN_INFO
        info_list = [self.G_INF, self.LRT_INF]
        inf_counter = 0
        lrt_types = [0, 0]
        info_types = ["CAP", "NLSL", "ELWL", "IMG5", "IMG6", "GEONS", "GEOEW", "LOOP", "LRT_Orig", "LRT_Orig",
                      "INF"]
        while inf_counter < len(info_types):
            if inf_counter == 8 or inf_counter == 9:
                curr_inf = info_list[1]
                lrt_types[inf_counter - 8] = phaser_gen_info_list[inf_counter]
                data_to_push = lrt_types
            else:
                curr_inf = info_list[0]
                data_to_push = phaser_gen_info_list[inf_counter]
            setattr(curr_inf, info_types[inf_counter], data_to_push)
            inf_counter += 1
        self.LRT_INF.lrt_orig_to_dir()

    def push_lrt_info(self):
        """the method uses the output LRT information of Phaser to push it into LRT_INF subclass and to each
        matching property in that subclass."""
        phaser_lrt_info_list = self.phsr_lst.LRT_INFO
        phaser_lrt_info_list.pop(0)
        lrt_info_types = ["CYC_TIME", "LRT_LOST_TIME", "LRT_HDWAY", "LRT_MCU", "GEN_LOST_TIME"]
        for lrt_inf in lrt_info_types:
            setattr(self.LRT_INF, lrt_inf, phaser_lrt_info_list[lrt_info_types.index(lrt_inf)])

    def push_street_names(self):
        """the method uses the output street names of Phaser to push it into phsr_lst subclass and to the
        matching property in that subclass (STREET)."""
        phaser_street_names_list = self.phsr_lst.STREETS
        dir_list = {"NO": self.NO, "SO": self.SO, "EA": self.EA, "WE": self.WE}
        dir_keys = list(dir_list.keys())
        for cur_dir in dir_keys:
            the_dir = dir_list[cur_dir]
            the_name = phaser_street_names_list[dir_keys.index(cur_dir)]
            setattr(the_dir, "NAME", the_name)

    def get_type_of_junc_for_choosing_slide(self):
        """This method checks about each direction in the junction whether it's empty or not (empty: no lanes or
        volumes); It creates a string that is later translated into a type of junction ＋,⊢,⊤,⊣,⊥ The method also
        checks for LRT in the junction (and it's direction), and for a metro around the junction. The method returns a
        number that matches the info about the junction and represents a matching slide in the diagrams template
        file. If the junction contains only three directions, the method updates a property that later will be
        used for the oneway function.
        """
        north = str(self.NO.empty_direction())
        south = str(self.SO.empty_direction())
        east = str(self.EA.empty_direction())
        west = str(self.WE.empty_direction())
        dir_exist = north + south + east + west
        types = {'1111': 1, '1110': 2, '0111': 3, '1101': 4, '1011': 5}
        if dir_exist not in types:
            # Fewer than 3 arms: add the first empty direction to reach a valid 3-arm config
            dir_list = list(dir_exist)
            for i in range(4):
                if dir_list[i] == '0':
                    dir_list[i] = '1'
                    candidate = ''.join(dir_list)
                    if candidate in types:
                        dir_exist = candidate
                        break
            self.WARNINGS.append(
                "אזהרה: הצומת הוגדר עם פחות משלוש זרועות. "
                "חישוב הדיאגרמה דורש לפחות 3 זרועות. "
                "זרוע וירטואלית נוספה אוטומטית לצורך החישוב — יש לוודא את נכונות הפלט."
            )
        if self.LRT_INF.LRT_Dir > 0:
            lrt_type = 1
        else:
            lrt_type = 0
        if self.LRT_INF.Metro_Dir > 0:
            metro_type = 2
        else:
            metro_type = 0
        junc_type = metro_type * 5 + lrt_type * 5 + types[dir_exist]
        if types[dir_exist] > 1:
            self.G_INF.ONEWAY = types[dir_exist]
        return junc_type

    def add_street_name_and_lrt(self, pres):
        """The method goes through all the shapes in Diagram pptx file and checks if it represents a direction
        name; It adds the matching name to each street, based on the name property of each direction. If the junction
        has lrt, the method adds the direction of the lrt. If the junction has only 3 direction, the method calls
        'is_oneway' method, to check if it's a oneway street
        """

        street_placeholders = {"NORTH_NAME": self.NO.NAME, "SOUTH_NAME": self.SO.NAME, "EAST_NAME": self.EA.NAME,
                               "WEST_NAME": self.WE.NAME, "RAKAL": self.LRT_INF.LRT_Dir}
        lrt_type_to_string = {
            0:  "",
            1:  "צפון ⇋ דרום",
            2:  "מזרח ⇋ מערב",
            3:  "צפ ⇋ דר, מז ⇋ מע",
            4:  "צפ ⇋ דר + כניסה ממזרח",
            5:  "צפ ⇋ דר + כניסה ממערב",
            6:  "מז ⇋ מע + כניסה מצפון",
            7:  "מז ⇋ מע + כניסה מדרום",
            8:  "צפון ⇋ מזרח",
            9:  "מזרח ⇋ דרום",
            10: "דרום ⇋ מערב",
            11: "מערב ⇋ צפון",
        }
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.name in street_placeholders.keys():
                    text_frame = shape.text_frame
                    text_frame.clear()
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    if shape.name == "RAKAL":
                        text_frame = shape.text_frame
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        font = run.font
                        font.bold = False
                        font.size = Pt(18)
                        font.color.rgb = RGBColor(228, 223, 211)
                        font.language_id = MSO_LANGUAGE_ID.HEBREW
                        text_frame.text_wrap = True
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        text_frame.text_wrap = True
                        font.name = 'Assistant'
                        run.text = str(lrt_type_to_string[street_placeholders[shape.name]])
                    else:
                        text_frame = shape.text_frame
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        font = run.font
                        font.bold = True
                        font.size = Pt(24)
                        font.language_id = MSO_LANGUAGE_ID.HEBREW
                        font.color.rgb = RGBColor(89, 89, 89)
                        font.name = 'Assistant'
                        if street_placeholders[shape.name] == 0:
                            run.text = ""
                        else:
                            run.text = str(street_placeholders[shape.name])
                        text_frame.text_wrap = True
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if self.G_INF.ONEWAY > 1:
            self.is_oneway(pres)
        pres.save("Street_Diagram.pptx")

    def add_morning_volumes(self, pres):
        """
        The method goes through all the shapes in Diagram pptx file and checks if it represents a volume of the
        morning counts; It adds the matching volume to each direction, based on the route property of each direction(
        L,T,R).
        """

        volume_placeholders = {"NORTH_MOR_R": self.NO.MOR.R, "NORTH_MOR_T": self.NO.MOR.T, "NORTH_MOR_L": self.NO.MOR.L,
                               "SOUTH_MOR_R": self.SO.MOR.R, "SOUTH_MOR_T": self.SO.MOR.T, "SOUTH_MOR_L": self.SO.MOR.L,
                               "EAST_MOR_R": self.EA.MOR.R, "EAST_MOR_T": self.EA.MOR.T, "EAST_MOR_L": self.EA.MOR.L,
                               "WEST_MOR_R": self.WE.MOR.R, "WEST_MOR_T": self.WE.MOR.T, "WEST_MOR_L": self.WE.MOR.L
                               }
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.name in volume_placeholders.keys():
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    font = run.font
                    font.language_id = MSO_LANGUAGE_ID.HEBREW
                    font.bold = True
                    font.italic = None
                    font.size = Pt(16)
                    font.color.rgb = RGBColor(255, 192, 0)
                    font.name = "Assistant"
                    if volume_placeholders[shape.name] == 0:
                        data_to_push = "-"
                    else:
                        data_to_push = str(volume_placeholders[shape.name])
                    run.text = data_to_push
        pres.save("Morn_Diagram.pptx")

    def add_evening_volumes(self, pres):
        """
        The method goes through all the shapes in Diagram pptx file and checks if it represents a volume of the
        evening counts; It adds the matching volume to each direction, based on the route property of each direction(
        L,T,R).
        """
        volume_placeholders = {"NORTH_EVE_R": self.NO.EVE.R, "NORTH_EVE_T": self.NO.EVE.T, "NORTH_EVE_L": self.NO.EVE.L,
                               "SOUTH_EVE_R": self.SO.EVE.R, "SOUTH_EVE_T": self.SO.EVE.T, "SOUTH_EVE_L": self.SO.EVE.L,
                               "EAST_EVE_R": self.EA.EVE.R, "EAST_EVE_T": self.EA.EVE.T, "EAST_EVE_L": self.EA.EVE.L,
                               "WEST_EVE_R": self.WE.EVE.R, "WEST_EVE_T": self.WE.EVE.T, "WEST_EVE_L": self.WE.EVE.L}
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.name in volume_placeholders.keys():
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    font = run.font
                    font.bold = True
                    font.italic = None
                    font.size = Pt(16)
                    font.color.rgb = RGBColor(200, 214, 223)
                    font.name = "Assistant"
                    if volume_placeholders[shape.name] == 0:
                        data_to_push = "-"
                    else:
                        data_to_push = str(volume_placeholders[shape.name])
                    run.text = data_to_push
        pres.save("Eve_Diagram.pptx")

    def add_direction_arrows(self, pres):
        """
        The method goes through all the shapes in Diagram pptx file and checks if it represents a string of arrows;
        It adds the matching arrow string after organizing them in the right logical order,
        using Organize_arrows_order method.
        """
        match_colors_to_type = {"White": RGBColor(255, 255, 255), "Yellow": RGBColor(250, 201, 49)}
        arrows_placeholders = {"NORTH_ARROWS": self.NO.LAN.Organize_arrows_order(),
                               "SOUTH_ARROWS": self.SO.LAN.Organize_arrows_order(),
                               "EAST_ARROWS": self.EA.LAN.Organize_arrows_order(),
                               "WEST_ARROWS": self.WE.LAN.Organize_arrows_order()}
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.name in arrows_placeholders.keys():
                    text_frame = shape.text_frame
                    text_frame.clear()
                    arrows_list = arrows_placeholders[shape.name]
                    for arrow in arrows_list:
                        text_frame = shape.text_frame
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        font = run.font
                        font.bold = False
                        font.size = Pt(50)
                        font.color.rgb = match_colors_to_type[arrow[2]]
                        font.name = 'Traffic Arrows 2 Med normal'
                        run.text = arrow[0] * arrow[1]
        pres.save("Dirc_Diagram.pptx")

    def is_oneway(self, pres):
        """the method is being called when the junction has three directions; It checks if the fourth direction is a
        oneway direction. If it is, it adds a matching oneway road to the final diagram.  """
        opt_oneway = self.G_INF.ONEWAY
        src = os.path.join(_JUNC_DIR, "Oneway_template") + os.sep
        type_dict = {2: "NoWest", 3: "NoNorth", 4: "NoEast", 5: "NoSouth"}
        OneWay = ""
        if type_dict[opt_oneway] == "NoNorth":
            if (self.EA.MOR.R > 0) or (self.EA.EVE.R > 0) \
                    or (self.WE.MOR.L > 0) or (self.WE.EVE.L > 0) \
                    or (self.SO.MOR.T > 0) or (self.SO.EVE.T > 0):
                OneWay = "North"

        if type_dict[opt_oneway] == "NoSouth":
            if (self.WE.MOR.R > 0) or (self.WE.EVE.R > 0) \
                    or (self.EA.MOR.L > 0) or (self.EA.EVE.L > 0) \
                    or (self.NO.MOR.T > 0) or (self.NO.EVE.T > 0):
                OneWay = "South"

        if type_dict[opt_oneway] == "NoEast":
            if (self.SO.MOR.R > 0) or (self.SO.EVE.R > 0) \
                    or (self.NO.MOR.L > 0) or (self.NO.EVE.L > 0) \
                    or (self.WE.MOR.T > 0) or (self.WE.EVE.T > 0):
                OneWay = "East"

        if type_dict[opt_oneway] == "NoWest":
            if (self.NO.MOR.R > 0) or (self.NO.EVE.R > 0) \
                    or (self.SO.MOR.L > 0) or (self.SO.EVE.L > 0) \
                    or (self.EA.MOR.T > 0) or (self.EA.EVE.T > 0):
                OneWay = "West"
        if OneWay != "":
            for slide in pres.slides:
                oneway_prop = {"North": [2915380, 0],  # North: [(slide.width - pic.width)/2, 0]
                               "South": [2915380, 4635945],
                               # South: [(slide.width - pic.width)/2, (slide.height - pic.height)]
                               "East": [4635945, 2915380],
                               # East: [(slide.width - pic.width), (slide.width - pic.width)/2]
                               "West": [0, 2915380]}  # West: [0, (slide.width - pic.width)/2]

                img_path = src + OneWay + "_one_way.png"
                slide.shapes.add_picture(img_path, oneway_prop[OneWay][0], oneway_prop[OneWay][1])


def run_diagram_pipeline(phsr_list, excel_properties):
    """Build the Diagram PPTX. Returns (junc_diagram, pptx_bytes)."""
    junc_diagram = Diagram(phsr_list, excel_properties)
    junc_diagram.push_arr()
    junc_diagram.push_vol()
    junc_diagram.push_general_info()
    junc_diagram.push_lrt_info()
    junc_diagram.push_street_names()
    try:
        create_new_diagram_template_file()
        prs = Presentation("Diagram_new_template.pptx")
        del_slides(prs, junc_diagram.get_type_of_junc_for_choosing_slide())
        prs = Presentation("Del_Diagram.pptx")
        junc_diagram.add_street_name_and_lrt(prs)
        prs = Presentation("Street_Diagram.pptx")
        junc_diagram.add_morning_volumes(prs)
        prs = Presentation("Morn_Diagram.pptx")
        junc_diagram.add_evening_volumes(prs)
        prs = Presentation("Eve_Diagram.pptx")
        junc_diagram.add_direction_arrows(prs)
        prs = Presentation("Dirc_Diagram.pptx")
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return junc_diagram, buf.getvalue()
    finally:
        delete_temp_diagram_pres()
