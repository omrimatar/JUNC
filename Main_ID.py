import time
from io import BytesIO
from Building_ID import *
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.dml.color import RGBColor
from pptx.util import Pt


class ID:
    """
      A class used to represent an ID for the created JUNC "
    """

    def __init__(self, phsr_list):
        self.__Project_name = ""
        self.__Project_number = ""
        self.__Author = ""
        self.__Date_time = str(time.strftime("%d/%m/%Y   %H:%M"))
        self.__Counter = 0
        self.__More_info = ""
        self.__Street_names = ""
        self.__phaser_info_list = phsr_list

    @property
    def PROJ_NAME(self):
        """Get the name of the project"""
        return self.__Project_name

    @PROJ_NAME.setter
    def PROJ_NAME(self, name):
        """Set the name of the project"""
        self.__Project_name = name

    @property
    def PROJ_NUM(self):
        """Get the project number"""
        return self.__Project_number

    @PROJ_NUM.setter
    def PROJ_NUM(self, number):
        """Set the project number"""
        self.__Project_number = number

    @property
    def AUTHOR(self):
        """Get the name of the author, based on the info written in office"""
        return self.__Author

    @AUTHOR.setter
    def AUTHOR(self, author):
        """Set the name of the author | DO NOT USE UNLESS NEEDED!"""
        self.__Author = author

    @property
    def DATETIME(self):
        """Get the time of the created JUNC"""
        return self.__Date_time

    @DATETIME.setter
    def DATETIME(self, datetime):
        """Set the time of the created JUNC | DO NOT USE UNLESS NEEDED!"""
        self.__Date_time = datetime

    @property
    def COUNT(self):
        """Get the number of tries or versions for this JUNC"""
        return self.__Counter

    @COUNT.setter
    def COUNT(self, count):
        """Set  the number of tries or versions for this JUNC"""
        self.__Counter = count

    @property
    def INFO(self):
        """Get the additional info about the JUNC"""
        return self.__More_info

    @INFO.setter
    def INFO(self, info):
        """Set the additional info about the JUNC"""
        self.__More_info = info

    @property
    def STREETS(self):
        """Get the names of the streets for this JUNC"""
        return self.__Street_names

    @STREETS.setter
    def STREETS(self, streets):
        """Set the names of the streets for this JUNC | DO NOT USE UNLESS NEEDED!"""
        self.__Street_names = streets

    @property
    def phsrlst(self):
        """Get the phaser list info"""
        return self.__phaser_info_list

    @phsrlst.setter
    def phsrlst(self, value):
        """Set the phaser list info"""
        self.__phaser_info_list = value

    def push_id_info(self, diagram):
        """the method uses the output section info of Phaser and the excel properties of the volume_calculator to
        push it into the right property of the class. """
        self.PROJ_NAME = self.phsrlst[0]
        self.PROJ_NUM = self.phsrlst[1]
        self.COUNT = self.phsrlst[3]
        self.INFO = self.phsrlst[4]
        self.AUTHOR = diagram.xlprop[0].lastModifiedBy
        streets = str(diagram.WE.NAME) + " · " + str(diagram.EA.NAME) + " · " + str(diagram.SO.NAME) + " · " + str(
            diagram.NO.NAME)
        self.STREETS = streets

    def add_info(self, pres):
        """
        The method goes through all the shapes in ID pptx file and checks if it represents the information
        (if the info was typed into the volume_calculator file) and adds the matching info.
        """
        placeholders = {"PROJECT_NUMBER": self.PROJ_NUM, "PROJECT_NAME": self.PROJ_NAME, "COUNTER": self.COUNT,
                        "AUTHOR": self.AUTHOR, "DATE_TIME": self.DATETIME, "MORE_INFO": self.INFO,
                        "STREET_NAMES": self.STREETS}
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.name in placeholders.keys():
                    name_of_shape = shape.name
                    text_to_push = str(placeholders[name_of_shape])
                    if text_to_push == "0":
                        text_to_push = "-"
                    text_frame = shape.text_frame
                    text_frame.clear()
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    font = run.font
                    font.bold = False
                    font.italic = None
                    font.size = Pt(11)
                    font.color.rgb = RGBColor(0, 0, 0)
                    font.name = "Assistant"
                    run.text = text_to_push
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.word_wrap = True
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.language_id = MSO_LANGUAGE_ID.HEBREW
        pres.save("id_info.pptx")


def run_id_pipeline(junc_table, junc_diagram):
    """Build the ID PPTX. Returns pptx_bytes."""
    new_info = junc_table.phsrlst.ID_INFO
    junc_id = ID(new_info)
    junc_id.push_id_info(junc_diagram)
    try:
        create_new_id_templates_file()
        prs = Presentation("id_template.pptx")
        junc_id.add_info(prs)
        prs = Presentation("id_info.pptx")
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()
    finally:
        delete_temp_id_pres()
