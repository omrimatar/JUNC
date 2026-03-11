from io import BytesIO
from Building_Table import *
from Section import *
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


class Table:
    """
      A class used to represent the structure of the table with different fields"
    """

    def __init__(self, phsr_list):
        """ The constructor of the Table class,called when a new instance of a class is created.
         To initialize, it needs the output of Phaser"""
        self.__Morning = Section()  # A property representing the Morning section in the table.
        self.__Evening = Section()  # A property representing the Evening section in the table.
        self.__LRT_status = False  # A property representing the status of the LRT. Affects the type of the table.
        self.__num_of_images = 3  # A property representing the number of images in the table.
        self.__phaser_info_list = phsr_list

    @property
    def MOR(self):
        """Get the info for the morning section"""
        return self.__Morning

    @MOR.setter
    def MOR(self, value):
        """Set the info for the morning section"""
        self.__Morning = value

    @property
    def EVE(self):
        """Get the info for the morning section"""
        return self.__Evening

    @EVE.setter
    def EVE(self, value):
        """Set info for the morning section"""
        self.__Evening = value

    @property
    def IS_LRT(self):
        """Get the info about the LRT status"""
        return self.__LRT_status

    @IS_LRT.setter
    def IS_LRT(self, value):
        """Set the info about the LRT status"""
        self.__LRT_status = value

    @property
    def IMG(self):
        """Get the number of the images in the table"""
        return self.__num_of_images

    @IMG.setter
    def IMG(self, value):
        """Set the number of the images in the table"""
        self.__num_of_images = value

    @property
    def phsrlst(self):
        """Get the phaser list info"""
        return self.__phaser_info_list

    @phsrlst.setter
    def phsrlst(self, value):
        """Set the phaser list info"""
        self.__phaser_info_list = value

    def push_deter_vol(self):
        """the method uses the output determining volumes of Phaser to push them into the right subclass of each images,
        divided to the morning and evening sections."""
        vol_img = ['A', 'B', 'C', 'D', 'E', 'F']
        time_section = {"MOR": self.phsrlst.MOR_DETER_VOL, "EVE": self.phsrlst.EVE_DETER_VOL}
        for cur_time in time_section.keys():
            for img in vol_img:
                cur_vol_img = "image" + img
                cur_section = getattr(self, cur_time)
                cur_img = getattr(cur_section, img)
                value_to_push = int(time_section[cur_time][cur_vol_img])
                setattr(cur_img, "VOL", value_to_push)

    def push_section_info(self):
        """the method uses the output section info of Phaser to push it into the right subclass of each section->
        divided to the morning and evening sections."""
        info_index = ["MOR_VOC", "MOR_TOT", "MOR_LRT", "EVE_VOC", "EVE_TOT", "EVE_LRT"]
        time_section = ["MOR", "EVE"]
        sections_info_list = ["VOC", "TOT", "LRT"]
        inf_count = 0
        for cur_time in time_section:
            for sect in sections_info_list:
                attr_time = getattr(self, cur_time)
                value_to_push = getattr(self.phsrlst, info_index[inf_count])
                if sect == "TOT" or sect == "LRT":
                    value_to_push = int(value_to_push)
                setattr(attr_time, sect, value_to_push)
                inf_count += 1
        self.MOR.set_los()
        self.EVE.set_los()

    def push_arrow_imgs(self):
        """the method uses the output arrows of Phaser to push them into the right subclass of each section,
        divided to morning and evening. """
        arrow_imgs = {"MOR": self.phsrlst.MOR_ARROW_TABLE, "EVE": self.phsrlst.EVE_ARROW_TABLE}
        for cur_time in arrow_imgs.keys():
            cur_section = getattr(self, cur_time)
            cur_section.split_img(arrow_imgs[cur_time])

    def deter_num_of_img(self):
        """The method checks how many images the output contains and sets the amount to self.IMG"""
        img_counter = 0

        vol_img = ['A', 'B', 'C', 'D', 'E', 'F']
        time_section = {"MOR": self.phsrlst.MOR_DETER_VOL, "EVE": self.phsrlst.EVE_DETER_VOL}
        for cur_time in time_section.keys():
            temp_count = 0
            for img in vol_img:
                cur_section = getattr(self, cur_time)
                cur_img = getattr(cur_section, img)
                value_to_check = getattr(cur_img, "VOL")
                if value_to_check > 0:
                    temp_count += 1
            if temp_count > img_counter:
                img_counter = temp_count
        self.IMG = img_counter

    def set_lrt_status(self, junc_diagram):
        """Sets the IS_LRT property, based on the info about it in Diagram"""
        if junc_diagram.LRT_INF.LRT_Dir > 0:
            self.IS_LRT = True
        else:
            self.IS_LRT = False

    def get_type_of_table_for_choosing_slide(self, junc_diagram):
        """The method chooses the type of table based on the info about num of images and LRT status"""
        self.deter_num_of_img()
        self.set_lrt_status(junc_diagram)
        table_slide_from_images = self.IMG - 1
        if self.IS_LRT:
            lrt_type = 1
        else:
            lrt_type = 0
        table_type = table_slide_from_images + lrt_type * 5
        return table_type

    def add_deter_volumes(self, pres):
        """
        The method goes through all the shapes in Table pptx file and checks if it represents a determining volume of
        the morning or evening in one of the images; It adds the matching volume to each image.
        """
        volume_placeholders = {"MOR_VOL_A": self.MOR.A.VOL, "MOR_VOL_B": self.MOR.B.VOL, "MOR_VOL_C": self.MOR.C.VOL,
                               "MOR_VOL_D": self.MOR.D.VOL, "MOR_VOL_E": self.MOR.E.VOL, "MOR_VOL_F": self.MOR.F.VOL,
                               "EVE_VOL_A": self.EVE.A.VOL, "EVE_VOL_B": self.EVE.B.VOL, "EVE_VOL_C": self.EVE.C.VOL,
                               "EVE_VOL_D": self.EVE.D.VOL, "EVE_VOL_E": self.EVE.E.VOL, "EVE_VOL_F": self.EVE.F.VOL
                               }
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.name in volume_placeholders.keys():
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    font = run.font
                    font.bold = False
                    font.italic = None
                    font.size = Pt(20)
                    font.color.rgb = RGBColor(0, 0, 0)
                    font.name = "Assistant"
                    run.text = str(volume_placeholders[shape.name])
        pres.save("Vol_Table.pptx")

    def add_table_info(self, pres):
        """
        The method goes through all the shapes in Table pptx file and checks if it represents the general information
        about the morning or evening sections; It adds the matching info.
        """
        info_placeholders = {"MOR_TOT": self.MOR.TOT, "MOR_VOC": self.MOR.VOC, "MOR_LOS": self.MOR.LOS,
                             "MOR_LRT": self.MOR.LRT, "EVE_TOT": self.EVE.TOT, "EVE_VOC": self.EVE.VOC,
                             "EVE_LOS": self.EVE.LOS, "EVE_LRT": self.EVE.LRT
                             }
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.name in info_placeholders.keys():
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    font = run.font
                    font.bold = False
                    font.italic = None
                    font.size = Pt(20)
                    font.color.rgb = RGBColor(0, 0, 0)
                    font.name = "Assistant"
                    run.text = str(info_placeholders[shape.name])
        pres.save("Info_Table.pptx")

    def add_table_arrows(self, JUNC_diagram, pres):
        """
        The method goes through all the shapes in Table pptx file and checks if it represents the string of arrows
        for one of the images. Before adding the arrows, it organizes the arrows in the logical order using
        organize_arrows_order_for_table function.
        """
        time_self = [self.MOR, self.EVE]
        for chosen_time in time_self:
            i = 0
            imgs_for_time = [chosen_time.A, chosen_time.B, chosen_time.C, chosen_time.D, chosen_time.E, chosen_time.F]
            while i < self.IMG:
                imgs_for_time[i].organize_arrows_order_for_table(JUNC_diagram, imgs_for_time[i])
                i += 1
        match_colors_to_type = {"White": RGBColor(0, 0, 0), "Yellow": RGBColor(250, 201, 49)}
        img_list = ["A", "B", "C", "D", "E", "F"]
        dir_list = ["NO", "SO", "EA", "WE"]
        time_list = ["MOR", "EVE"]

        for slide in pres.slides:
            for shape in slide.shapes:
                if (str(shape.name)[:3]) in time_list:
                    if (str(shape.name)[-1]) in img_list:
                        if str(shape.name)[4:6] in dir_list:
                            time_attr = getattr(self, str(shape.name)[:3])
                            img_attr = getattr(time_attr, str(shape.name)[-1])
                            dirc_attr = getattr(img_attr, str(shape.name)[4:6])
                            if dirc_attr:
                                text_frame = shape.text_frame
                                text_frame.clear()
                                arrows_list = dirc_attr
                                for arrow in arrows_list:
                                    text_frame = shape.text_frame
                                    p = text_frame.paragraphs[0]
                                    run = p.add_run()
                                    font = run.font
                                    font.bold = False
                                    font.size = Pt(24)
                                    font.color.rgb = match_colors_to_type[arrow[2]]
                                    font.name = 'Traffic Arrows 2 Med normal'
                                    run.text = arrow[0] * arrow[1]
        pres.save("Dirc_Table.pptx")


def run_table_pipeline(junc_diagram):
    """Build the Table PPTX. Returns (junc_table, pptx_bytes)."""
    junc_table = Table(junc_diagram.phsr_lst)
    junc_table.push_deter_vol()
    junc_table.push_section_info()
    junc_table.push_arrow_imgs()
    try:
        create_new_table_templates_file()
        prs = Presentation("Table_new_template.pptx")
        del_slides_table(prs, junc_table.get_type_of_table_for_choosing_slide(junc_diagram))
        prs = Presentation("Del_Table.pptx")
        junc_table.add_deter_volumes(prs)
        prs = Presentation("Vol_Table.pptx")
        junc_table.add_table_info(prs)
        prs = Presentation("Info_Table.pptx")
        junc_table.add_table_arrows(junc_diagram, prs)
        prs = Presentation("Dirc_Table.pptx")
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return junc_table, buf.getvalue()
    finally:
        delete_temp_table_pres()
